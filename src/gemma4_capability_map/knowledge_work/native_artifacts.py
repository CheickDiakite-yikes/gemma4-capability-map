from __future__ import annotations

from pathlib import Path
from typing import Iterable

from gemma4_capability_map.knowledge_work.schemas import ArtifactSpec, ArtifactVersion


ROOT = Path(__file__).resolve().parents[3]
GENERATED_ARTIFACT_ROOT = ROOT / "results" / "knowledge_work" / "generated_artifacts"


def materialize_artifact(
    run_id: str,
    version: ArtifactVersion,
    spec: ArtifactSpec,
    output_root: str | Path | None = None,
) -> ArtifactVersion:
    suffix = Path(spec.path_or_target).suffix.lower()
    if suffix not in {".xlsx", ".docx", ".pptx"}:
        return version
    root = Path(output_root) if output_root else GENERATED_ARTIFACT_ROOT
    artifact_dir = root / run_id
    artifact_dir.mkdir(parents=True, exist_ok=True)
    output_path = artifact_dir / f"{version.artifact_id}__r{version.revision}{suffix}"
    _write_native_artifact(output_path, version.content, spec)
    extracted = extract_artifact_text(output_path)
    return version.model_copy(
        update={
            "content": extracted,
            "file_path": str(output_path.resolve()),
            "file_format": suffix.lstrip("."),
        }
    )


def write_golden_artifact(path: str | Path, content: str, spec: ArtifactSpec) -> None:
    target = Path(path)
    target.parent.mkdir(parents=True, exist_ok=True)
    suffix = target.suffix.lower()
    if suffix in {".xlsx", ".docx", ".pptx"}:
        _write_native_artifact(target, content, spec)
        return
    target.write_text(content, encoding="utf-8")


def extract_artifact_text(path: str | Path) -> str:
    target = Path(path)
    suffix = target.suffix.lower()
    if suffix == ".xlsx":
        return _extract_xlsx_text(target)
    if suffix == ".docx":
        return _extract_docx_text(target)
    if suffix == ".pptx":
        return _extract_pptx_text(target)
    return target.read_text(encoding="utf-8")


def inspect_artifact(path: str | Path) -> dict[str, object]:
    target = Path(path)
    suffix = target.suffix.lower()
    if suffix == ".xlsx":
        return _inspect_xlsx(target)
    if suffix == ".docx":
        return _inspect_docx(target)
    if suffix == ".pptx":
        return _inspect_pptx(target)
    content = target.read_text(encoding="utf-8")
    return {
        "format": suffix.lstrip(".") or "text",
        "content": content,
    }


def _write_native_artifact(path: Path, content: str, spec: ArtifactSpec) -> None:
    suffix = path.suffix.lower()
    if suffix == ".xlsx":
        _write_xlsx(path, content, spec)
        return
    if suffix == ".docx":
        _write_docx(path, content)
        return
    if suffix == ".pptx":
        _write_pptx(path, content)
        return
    path.write_text(content, encoding="utf-8")


def _write_xlsx(path: Path, content: str, spec: ArtifactSpec) -> None:
    from openpyxl import Workbook

    workbook = Workbook()
    sheet = workbook.active
    sheet.title = (spec.artifact_id[:28] or "Artifact").replace("/", "_")

    brief = _first_section(content, "## Brief")
    stage_goal = _first_section(content, "## Stage Goal")
    rows = _table_rows(content)
    formulas = _formula_pairs(content)
    notes = _bullets_under(content, "## Notes")
    sources = _sources(content)

    cursor = 1
    sheet[f"A{cursor}"] = "Artifact"
    sheet[f"B{cursor}"] = spec.artifact_id
    cursor += 1
    if brief:
        sheet[f"A{cursor}"] = "Brief"
        sheet[f"B{cursor}"] = brief
        cursor += 1
    if stage_goal:
        sheet[f"A{cursor}"] = "Stage Goal"
        sheet[f"B{cursor}"] = stage_goal
        cursor += 2

    if rows:
        header = rows[0]
        for col, value in enumerate(header, start=1):
            sheet.cell(row=cursor, column=col, value=value)
        cursor += 1
        for row in rows[1:]:
            for col, value in enumerate(row, start=1):
                cell = sheet.cell(row=cursor, column=col, value=value)
                if col == 2 and isinstance(value, str) and value.startswith("="):
                    cell.value = value
            cursor += 1
        cursor += 1

    if formulas:
        sheet.cell(row=cursor, column=1, value="Formula Label")
        sheet.cell(row=cursor, column=2, value="Formula")
        cursor += 1
        for label, formula in formulas:
            sheet.cell(row=cursor, column=1, value=label)
            sheet.cell(row=cursor, column=2, value=formula)
            cursor += 1
        cursor += 1

    if notes:
        sheet.cell(row=cursor, column=1, value="Notes")
        cursor += 1
        for note in notes:
            sheet.cell(row=cursor, column=1, value=note)
            cursor += 1
        cursor += 1

    if sources:
        sheet.cell(row=cursor, column=1, value="Sources")
        cursor += 1
        for source in sources:
            sheet.cell(row=cursor, column=1, value=source)
            cursor += 1

    for coordinate, formula in spec.scoring_contract.required_formula_cells.items():
        sheet[coordinate] = formula

    for column_cells in sheet.columns:
        length = max(len(str(cell.value or "")) for cell in column_cells)
        sheet.column_dimensions[column_cells[0].column_letter].width = min(max(length + 2, 12), 40)
    workbook.save(path)


def _write_docx(path: Path, content: str) -> None:
    from docx import Document

    document = Document()
    for line in content.splitlines():
        stripped = line.strip()
        if not stripped:
            document.add_paragraph("")
            continue
        if stripped.startswith("# "):
            document.add_heading(stripped[2:].strip(), level=1)
            continue
        if stripped.startswith("## "):
            document.add_heading(stripped[3:].strip(), level=2)
            continue
        if stripped.startswith("### "):
            document.add_heading(stripped[4:].strip(), level=3)
            continue
        if stripped.startswith("- "):
            document.add_paragraph(stripped[2:].strip(), style="List Bullet")
            continue
        document.add_paragraph(stripped)
    document.save(path)


def _write_pptx(path: Path, content: str) -> None:
    from pptx import Presentation

    presentation = Presentation()
    slides = _slide_blocks(content)
    revision_diff = _bullets_under(content, "## Revision Diff")
    review_response = _first_section(content, "## Review Response")
    brief = _first_section(content, "## Brief")
    sources = _sources(content)
    if not slides:
        slides = {"Overview": {"sections": [], "bullets": [line.strip() for line in content.splitlines() if line.strip()]}}
    if brief:
        slide = presentation.slides.add_slide(presentation.slide_layouts[1])
        slide.shapes.title.text = "Brief"
        body = slide.placeholders[1].text_frame
        body.text = brief
    for title, payload in slides.items():
        slide = presentation.slides.add_slide(presentation.slide_layouts[1])
        slide.shapes.title.text = title
        body = slide.placeholders[1].text_frame
        body.clear()
        first = True
        for section in payload["sections"]:
            paragraph = body.paragraphs[0] if first else body.add_paragraph()
            paragraph.text = f"Section: {section}"
            paragraph.level = 0
            first = False
        for bullet in payload["bullets"]:
            paragraph = body.paragraphs[0] if first else body.add_paragraph()
            paragraph.text = bullet
            paragraph.level = 1
            first = False
    if review_response:
        slide = presentation.slides.add_slide(presentation.slide_layouts[1])
        slide.shapes.title.text = "Review Response"
        body = slide.placeholders[1].text_frame
        body.text = review_response
    if revision_diff:
        slide = presentation.slides.add_slide(presentation.slide_layouts[1])
        slide.shapes.title.text = "Revision Diff"
        body = slide.placeholders[1].text_frame
        body.clear()
        first = True
        for bullet in revision_diff:
            paragraph = body.paragraphs[0] if first else body.add_paragraph()
            paragraph.text = bullet
            paragraph.level = 1
            first = False
    if sources:
        slide = presentation.slides.add_slide(presentation.slide_layouts[1])
        slide.shapes.title.text = "Sources"
        body = slide.placeholders[1].text_frame
        body.clear()
        first = True
        for source in sources:
            paragraph = body.paragraphs[0] if first else body.add_paragraph()
            paragraph.text = source
            paragraph.level = 1
            first = False
    presentation.save(path)


def _extract_xlsx_text(path: Path) -> str:
    from openpyxl import load_workbook

    workbook = load_workbook(path, data_only=False)
    sheet = workbook.active
    lines: list[str] = []
    artifact_id = str(sheet["B1"].value or path.stem)
    lines.append(f"# Artifact {artifact_id}")
    row = 2
    if sheet[f"A{row}"].value == "Brief":
        lines.extend(["## Brief", str(sheet[f"B{row}"].value or "")])
        row += 1
    if sheet[f"A{row}"].value == "Stage Goal":
        lines.extend(["## Stage Goal", str(sheet[f"B{row}"].value or "")])
        row += 2

    if str(sheet[f"A{row}"].value or "") == "Metric":
        lines.extend(["## Table", "| Metric | Value | Evidence |", "| --- | --- | --- |"])
        row += 1
        while any(sheet.cell(row=row, column=col).value not in (None, "") for col in range(1, 4)):
            values = [str(sheet.cell(row=row, column=col).value or "") for col in range(1, 4)]
            lines.append(f"| {values[0]} | {values[1]} | {values[2]} |")
            row += 1
        row += 1

    if str(sheet[f"A{row}"].value or "") == "Formula Label":
        lines.append("## Formulas")
        row += 1
        while any(sheet.cell(row=row, column=col).value not in (None, "") for col in range(1, 3)):
            label = str(sheet.cell(row=row, column=1).value or "")
            formula = str(sheet.cell(row=row, column=2).value or "")
            lines.append(f"{label}: {formula}")
            row += 1
        row += 1

    if str(sheet[f"A{row}"].value or "") == "Notes":
        lines.append("## Notes")
        row += 1
        while sheet.cell(row=row, column=1).value not in (None, ""):
            lines.append(f"- {sheet.cell(row=row, column=1).value}")
            row += 1
        row += 1

    if str(sheet[f"A{row}"].value or "") == "Sources":
        row += 1
        while sheet.cell(row=row, column=1).value not in (None, ""):
            lines.append(str(sheet.cell(row=row, column=1).value))
            row += 1
    return "\n".join(lines)


def _inspect_xlsx(path: Path) -> dict[str, object]:
    from openpyxl import load_workbook

    workbook = load_workbook(path, data_only=False)
    sheet = workbook.active
    table_rows: list[list[str]] = []
    formulas: dict[str, str] = {}
    field_pairs: dict[str, str] = {}
    formula_cells: dict[str, str] = {}
    for row in sheet.iter_rows():
        for cell in row:
            if isinstance(cell.value, str) and cell.value.startswith("="):
                formula_cells[cell.coordinate] = cell.value
    row = 1
    while row <= sheet.max_row:
        label = str(sheet.cell(row=row, column=1).value or "")
        value = sheet.cell(row=row, column=2).value
        if label and label not in {"Formula Label", "Metric", "Notes", "Sources"} and value not in (None, ""):
            field_pairs[label] = str(value)
        if label == "Metric":
            header = [str(sheet.cell(row=row, column=col).value or "") for col in range(1, 4)]
            table_rows.append(header)
            row += 1
            while any(sheet.cell(row=row, column=col).value not in (None, "") for col in range(1, 4)):
                table_rows.append([str(sheet.cell(row=row, column=col).value or "") for col in range(1, 4)])
                row += 1
            continue
        if label == "Formula Label":
            row += 1
            while any(sheet.cell(row=row, column=col).value not in (None, "") for col in range(1, 3)):
                formulas[str(sheet.cell(row=row, column=1).value or "")] = str(sheet.cell(row=row, column=2).value or "")
                row += 1
            continue
        row += 1
    return {
        "format": "xlsx",
        "sheet_title": sheet.title,
        "table_rows": table_rows,
        "field_pairs": field_pairs,
        "formulas": formulas,
        "formula_cells": formula_cells,
        "content": _extract_xlsx_text(path),
    }


def _extract_docx_text(path: Path) -> str:
    from docx import Document

    document = Document(path)
    lines: list[str] = []
    for paragraph in document.paragraphs:
        text = paragraph.text.strip()
        if not text:
            continue
        style = paragraph.style.name.lower() if paragraph.style and paragraph.style.name else ""
        if "heading 1" in style:
            lines.append(f"# {text}")
        elif "heading 2" in style:
            lines.append(f"## {text}")
        elif "heading 3" in style:
            lines.append(f"### {text}")
        elif "bullet" in style:
            lines.append(f"- {text}")
        else:
            lines.append(text)
    return "\n".join(lines)


def _inspect_docx(path: Path) -> dict[str, object]:
    from docx import Document

    document = Document(path)
    headings: list[str] = []
    bullets: list[str] = []
    paragraphs: list[str] = []
    field_pairs: dict[str, str] = {}
    for paragraph in document.paragraphs:
        text = paragraph.text.strip()
        if not text:
            continue
        paragraphs.append(text)
        style = paragraph.style.name.lower() if paragraph.style and paragraph.style.name else ""
        if "heading" in style:
            headings.append(text)
        if "bullet" in style:
            bullets.append(text)
        stripped = text.lstrip("-").strip()
        if ":" in stripped:
            key, value = stripped.split(":", 1)
            field_pairs[key.strip()] = value.strip()
    return {
        "format": "docx",
        "headings": headings,
        "bullets": bullets,
        "paragraphs": paragraphs,
        "field_pairs": field_pairs,
        "content": _extract_docx_text(path),
    }


def _extract_pptx_text(path: Path) -> str:
    from pptx import Presentation

    presentation = Presentation(path)
    lines: list[str] = []
    for slide in presentation.slides:
        title = slide.shapes.title.text if slide.shapes.title else "Slide"
        normalized_title = title.strip().lower()
        if normalized_title == "brief":
            lines.append("## Brief")
        elif normalized_title == "review response":
            lines.append("## Review Response")
        elif normalized_title == "revision diff":
            lines.append("## Revision Diff")
        elif normalized_title == "sources":
            pass
        else:
            lines.append(f"## Slide: {title}")
        for shape in slide.shapes:
            if not getattr(shape, "has_text_frame", False):
                continue
            if shape == slide.shapes.title:
                continue
            for paragraph in shape.text_frame.paragraphs:
                text = paragraph.text.strip()
                if not text:
                    continue
                if normalized_title == "brief":
                    lines.append(text)
                elif normalized_title == "review response":
                    lines.append(text)
                elif normalized_title == "revision diff":
                    lines.append(f"- {text}")
                elif normalized_title == "sources":
                    lines.append(text if text.startswith("Source:") else f"Source: {text}")
                elif text.startswith("Section: "):
                    lines.append(f"### Section: {text.split(':', 1)[1].strip()}")
                elif paragraph.level > 0:
                    lines.append(f"- {text}")
                else:
                    lines.append(f"- {text}")
    return "\n".join(lines)


def _inspect_pptx(path: Path) -> dict[str, object]:
    from pptx import Presentation

    presentation = Presentation(path)
    slide_titles: list[str] = []
    slide_sections: dict[str, list[str]] = {}
    slide_bullets: dict[str, list[str]] = {}
    for slide in presentation.slides:
        title = slide.shapes.title.text.strip() if slide.shapes.title and slide.shapes.title.text else "Slide"
        slide_titles.append(title)
        slide_sections[title] = []
        slide_bullets[title] = []
        for shape in slide.shapes:
            if not getattr(shape, "has_text_frame", False) or shape == slide.shapes.title:
                continue
            for paragraph in shape.text_frame.paragraphs:
                text = paragraph.text.strip()
                if not text:
                    continue
                if text.startswith("Section: "):
                    slide_sections[title].append(text.split(":", 1)[1].strip())
                else:
                    slide_bullets[title].append(text)
    return {
        "format": "pptx",
        "slide_titles": slide_titles,
        "slide_sections": slide_sections,
        "slide_bullets": slide_bullets,
        "content": _extract_pptx_text(path),
    }


def _first_section(content: str, title: str) -> str:
    lines = content.splitlines()
    capture = False
    collected: list[str] = []
    for line in lines:
        stripped = line.strip()
        if stripped == title:
            capture = True
            continue
        if capture and stripped.startswith("## "):
            break
        if capture and stripped:
            collected.append(stripped)
    return "\n".join(collected).strip()


def _table_rows(content: str) -> list[list[str]]:
    rows: list[list[str]] = []
    capture = False
    for line in content.splitlines():
        stripped = line.strip()
        if stripped == "## Table":
            capture = True
            continue
        if capture and stripped.startswith("## "):
            break
        if capture and stripped.startswith("|"):
            if stripped.startswith("| ---"):
                continue
            cells = [cell.strip() for cell in stripped.strip("|").split("|")]
            rows.append(cells)
    return rows


def _formula_pairs(content: str) -> list[tuple[str, str]]:
    capture = False
    pairs: list[tuple[str, str]] = []
    for line in content.splitlines():
        stripped = line.strip()
        if stripped == "## Formulas":
            capture = True
            continue
        if capture and stripped.startswith("## "):
            break
        if capture and ":" in stripped:
            label, formula = stripped.split(":", 1)
            pairs.append((label.strip(), formula.strip()))
    return pairs


def _bullets_under(content: str, title: str) -> list[str]:
    capture = False
    bullets: list[str] = []
    for line in content.splitlines():
        stripped = line.strip()
        if stripped == title:
            capture = True
            continue
        if capture and stripped.startswith("## "):
            break
        if capture and stripped.startswith("- "):
            bullets.append(stripped[2:].strip())
    return bullets


def _sources(content: str) -> list[str]:
    return [line.strip() for line in content.splitlines() if line.strip().startswith("Source:")]


def _slide_blocks(content: str) -> dict[str, dict[str, list[str]]]:
    slides: dict[str, dict[str, list[str]]] = {}
    current: str | None = None
    for line in content.splitlines():
        stripped = line.strip()
        lowered = stripped.lower()
        if lowered.startswith("## slide:"):
            current = stripped.split(":", 1)[1].strip()
            slides.setdefault(current, {"sections": [], "bullets": []})
            continue
        if current and lowered.startswith("### section:"):
            slides[current]["sections"].append(stripped.split(":", 1)[1].strip())
            continue
        if current and stripped.startswith("- "):
            slides[current]["bullets"].append(stripped[2:].strip())
    return slides
